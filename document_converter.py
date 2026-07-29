#!/usr/bin/env python3
"""
Document Converter Tool

PPTX, DOCX, PDF 문서를 Markdown 으로 변환합니다.
페이지/슬라이드 번호 정보를 포함합니다.

Usage:
    python functions/tools/document_converter.py <file> [output_dir]
    python functions/tools/document_converter.py --batch <input_dir> [output_dir]
"""

from __future__ import annotations

import re
import sys
from pathlib import Path
from typing import Optional

try:
    from docx import Document
    from docx.oxml.ns import qn
except ImportError:
    print("Error: python-docx not installed. Run: pip install python-docx")
    sys.exit(1)


# ============================================================================
# DOCX → MD 변환
# ============================================================================

def extract_page_from_heading_text(text: str) -> tuple[str, Optional[int]]:
    """Extract page number from heading text if present."""
    match = re.match(r'^(.+?)\s+(\d+)$', text.strip())
    if match:
        title = match.group(1).strip()
        page_str = match.group(2)
        page_num = int(page_str)
        if 1 <= page_num <= 10000:
            return title, page_num
    return text, None


def convert_docx_to_md(docx_path: Path, include_page_numbers: bool = True) -> str:
    """
    Convert DOCX to Markdown with optional page number tracking.
    
    Args:
        docx_path: Path to DOCX file
        include_page_numbers: If True, add page numbers to headings (e.g., "## Title\t33")
    
    Returns:
        Markdown content
    """
    doc = Document(docx_path)
    md_lines: list[str] = []
    last_known_page = 1
    
    page_number_pattern = re.compile(r'^\s*(\d+)\s*$')
    
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        
        if page_number_pattern.match(text):
            continue
        
        style_name = para.style.name if para.style else ""
        
        # Heading detection
        heading_match = re.match(r"^Heading\s*(\d+)", style_name, re.IGNORECASE)
        if heading_match:
            level = int(heading_match.group(1))
            prefix = "#" * level
            if include_page_numbers:
                clean_text, page_num = extract_page_from_heading_text(text)
                if page_num:
                    last_known_page = page_num
                    md_lines.append(f"{prefix} {clean_text}\t{page_num}")
                else:
                    md_lines.append(f"{prefix} {clean_text}\t{last_known_page}")
            else:
                md_lines.append(f"{prefix} {text}")
            continue
        
        # 3GPP section numbers
        section_match = re.match(r"^(\d+(\.\d+)+)\s+(.+)$", text)
        if section_match and len(section_match.group(1).split(".")) <= 6:
            section_num = section_match.group(1)
            section_title = section_match.group(3)
            level = min(len(section_num.split(".")), 6)
            prefix = "#" * level
            if include_page_numbers:
                clean_title, page_num = extract_page_from_heading_text(section_title)
                if page_num:
                    last_known_page = page_num
                    md_lines.append(f"{prefix} {clean_title}\t{page_num}")
                else:
                    md_lines.append(f"{prefix} {clean_title}\t{last_known_page}")
            else:
                md_lines.append(f"{prefix} {section_title}")
            continue
        
        # Bold text headers
        if any(run.bold for run in para.runs if run.text.strip()):
            bold_text = "".join(run.text for run in para.runs if run.bold)
            if bold_text.strip() and len(bold_text.strip()) < 100:
                if include_page_numbers:
                    clean_text, page_num = extract_page_from_heading_text(text)
                    if page_num:
                        last_known_page = page_num
                        md_lines.append(f"## {clean_text}\t{page_num}")
                    else:
                        md_lines.append(f"## {clean_text}\t{last_known_page}")
                else:
                    md_lines.append(f"## {text}")
                continue
        
        # Regular paragraph
        md_lines.append(text)
    
    # Tables
    for table in doc.tables:
        md_lines.append("")
        md_lines.append(f"## Table\t{last_known_page}" if include_page_numbers else "## Table")
        md_lines.append(convert_table_to_md(table))
        md_lines.append("")
    
    return "\n".join(md_lines)


def convert_table_to_md(table) -> str:
    """Convert DOCX table to Markdown."""
    md_lines: list[str] = []
    rows = []
    for row in table.rows:
        cells = [cell.text.strip().replace("|", "\\|").replace("\n", " ") for cell in row.cells]
        rows.append(cells)
    
    if not rows:
        return ""
    
    max_cols = max(len(row) for row in rows)
    for row in rows:
        while len(row) < max_cols:
            row.append("")
    
    for i, row in enumerate(rows):
        md_lines.append("| " + " | ".join(row) + " |")
        if i == 0:
            md_lines.append("| " + " | ".join("---" for _ in row) + " |")
    
    return "\n".join(md_lines)


# ============================================================================
# PPTX → MD 변환 (read_pptx.py 와 통합)
# ============================================================================

def convert_pptx_to_md(pptx_path: Path) -> str:
    """
    Convert PPTX to Markdown with slide numbers.
    
    Args:
        pptx_path: Path to PPTX file
    
    Returns:
        Markdown content with slide numbers
    """
    try:
        from pptx import Presentation
    except ImportError:
        print("Error: python-pptx not installed. Run: pip install python-pptx")
        sys.exit(1)
    
    prs = Presentation(pptx_path)
    md_lines: list[str] = []
    
    for i, slide in enumerate(prs.slides, 1):
        md_lines.append(f"## Slide {i}")
        md_lines.append("")
        
        # Title
        if slide.shapes.title and slide.shapes.title.text:
            md_lines.append(f"### {slide.shapes.title.text}")
            md_lines.append("")
        
        # Content
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                md_lines.append(shape.text)
                md_lines.append("")
        
        md_lines.append("")
    
    return "\n".join(md_lines)


# ============================================================================
# PDF → MD 변환 (convert_pdf_to_md.py 와 통합)
# ============================================================================

def convert_pdf_to_md(pdf_path: Path) -> str:
    """
    Convert PDF to Markdown with page numbers.
    
    Args:
        pdf_path: Path to PDF file
    
    Returns:
        Markdown content with page numbers
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        print("Error: PyMuPDF not installed. Run: pip install PyMuPDF")
        sys.exit(1)
    
    doc = fitz.open(pdf_path)
    md_lines: list[str] = []
    
    for i, page in enumerate(doc, 1):
        md_lines.append(f"<!-- Page {i} -->")
        md_lines.append("")
        md_lines.append(page.get_text())
        md_lines.append("")
    
    doc.close()
    return "\n".join(md_lines)


# ============================================================================
# 메인 함수
# ============================================================================

def main():
    if len(sys.argv) < 2:
        print("Usage: python document_converter.py <file> [output_dir]")
        print("       python document_converter.py --batch <input_dir> [output_dir]")
        sys.exit(1)
    
    if sys.argv[1] == "--batch":
        if len(sys.argv) < 3:
            print("Error: input directory required for batch mode")
            sys.exit(1)
        
        input_dir = Path(sys.argv[2])
        output_dir = Path(sys.argv[3]) if len(sys.argv) > 3 else input_dir
        
        if not input_dir.exists():
            print(f"Error: Directory not found: {input_dir}")
            sys.exit(1)
        
        print(f"Batch converting files in: {input_dir}")
        print(f"Output directory: {output_dir}")
        print()
        
        converted = 0
        for file_path in input_dir.iterdir():
            if file_path.is_file():
                suffix = file_path.suffix.lower()
                if suffix == '.docx':
                    output_path = output_dir / f"{file_path.stem}.md"
                    md_content = convert_docx_to_md(file_path)
                    output_path.write_text(md_content, encoding='utf-8')
                    print(f"[OK] {file_path.name} -> {output_path.name}")
                    converted += 1
                elif suffix == '.pptx':
                    output_path = output_dir / f"{file_path.stem}.md"
                    md_content = convert_pptx_to_md(file_path)
                    output_path.write_text(md_content, encoding='utf-8')
                    print(f"[OK] {file_path.name} -> {output_path.name}")
                    converted += 1
                elif suffix == '.pdf':
                    output_path = output_dir / f"{file_path.stem}.md"
                    md_content = convert_pdf_to_md(file_path)
                    output_path.write_text(md_content, encoding='utf-8')
                    print(f"[OK] {file_path.name} -> {output_path.name}")
                    converted += 1
        
        print()
        print(f"Converted {converted} files")
    
    else:
        file_path = Path(sys.argv[1])
        output_dir = Path(sys.argv[2]) if len(sys.argv) > 2 else file_path.parent
        
        if not file_path.exists():
            print(f"Error: File not found: {file_path}")
            sys.exit(1)
        
        suffix = file_path.suffix.lower()
        if suffix == '.docx':
            output_path = output_dir / f"{file_path.stem}.md"
            md_content = convert_docx_to_md(file_path)
            output_path.write_text(md_content, encoding='utf-8')
            print(f"Converted {len(md_content)} characters")
            print(f"Output: {output_path}")
        elif suffix == '.pptx':
            output_path = output_dir / f"{file_path.stem}.md"
            md_content = convert_pptx_to_md(file_path)
            output_path.write_text(md_content, encoding='utf-8')
            print(f"Converted {len(md_content)} characters")
            print(f"Output: {output_path}")
        elif suffix == '.pdf':
            output_path = output_dir / f"{file_path.stem}.md"
            md_content = convert_pdf_to_md(file_path)
            output_path.write_text(md_content, encoding='utf-8')
            print(f"Converted {len(md_content)} characters")
            print(f"Output: {output_path}")
        else:
            print(f"Error: Unsupported file type: {suffix}")
            sys.exit(1)


if __name__ == "__main__":
    main()
