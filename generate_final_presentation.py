from __future__ import annotations
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 삼성 블루 컬러
SAMSUNG_BLUE = RGBColor(0, 24, 107)  # #00186B
SAMSUNG_LIGHT_BLUE = RGBColor(0, 82, 147)  # #005293
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(128, 128, 128)

def set_title_format(shape, text):
    """제목 서식 설정"""
    shape.text = text
    for paragraph in shape.text_frame.paragraphs:
        paragraph.font.bold = True
        paragraph.font.size = Pt(32)
        paragraph.font.color.rgb = SAMSUNG_BLUE
        paragraph.alignment = PP_ALIGN.CENTER

def set_subtitle_format(shape, text):
    """부제목 서식 설정"""
    shape.text = text
    for paragraph in shape.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = GRAY
        paragraph.alignment = PP_ALIGN.CENTER

def set_body_format(shape, text, font_size=14):
    """본문 서식 설정"""
    shape.text = text
    for paragraph in shape.text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = RGBColor(0, 0, 0)

def add_section_header(slide, section_num, section_title):
    """섹션 헤더 추가"""
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = SAMSUNG_BLUE
    shape.line.color.rgb = SAMSUNG_BLUE
    shape.text = f"{section_num}. {section_title}"
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.size = Pt(20)
    shape.text_frame.paragraphs[0].font.color.rgb = WHITE
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

def add_content_box(slide, left, top, width, height, title, content, title_color=SAMSUNG_BLUE):
    """콘텐츠 박스 추가"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = title_color
    shape.line.width = Pt(2)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = title_color
    
    for line in content.split('\n'):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_before = Pt(4)

def add_arrow(slide, left, top, width, height, direction='right'):
    """화살표 추가"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = SAMSUNG_LIGHT_BLUE
    shape.line.color.rgb = SAMSUNG_LIGHT_BLUE
    return shape

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

# ==================== Slide 1: 표지 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 슬라이드

# 제목
title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
set_title_format(title_box, "AI 기반 TC Review & Enhancement 시스템")

# 부제목
subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.333), Inches(1))
set_subtitle_format(subtitle_box, "3-Phase Block Workflow 와 품질 개선 알고리즘")

# 발표자 정보
presenter_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11.333), Inches(0.8))
set_subtitle_format(presenter_box, "Samsung Networks Business RVT-RAS | 김주영")

# ==================== Slide 2: 01. Background ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "01", "Background - TC 생성 자동화의 필요성")

# 3 가지 pain point
add_content_box(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(4), 
    "기능 범위 확대", 
    "• 새로운 rAPP 신규 개발\n• S/W 업데이트\n• 신규 요구사항 증가\n• 검증 대상 FR 증가",
    SAMSUNG_BLUE)

add_content_box(slide, Inches(4.6), Inches(1.5), Inches(3.8), Inches(4),
    "문서 해석 부담",
    "• HLD/DLD/Algorithm 분석 필요\n• 3GPP 규격과 용어 학습\n• 복합 Case 고려 필요",
    SAMSUNG_BLUE)

add_content_box(slide, Inches(8.7), Inches(1.5), Inches(3.8), Inches(4),
    "반복 작성 부담",
    "• 유사 TC 구조 반복\n• FR 1 개당 3~4 시간 소요\n• TC 누적 증가",
    SAMSUNG_BLUE)

# ==================== Slide 3: 02. Problem Situation ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "02", "Problem Situation - TC Lifecycle 병목")

# TC 생성 박스
shape1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2), Inches(3.5), Inches(2))
shape1.fill.solid()
shape1.fill.fore_color.rgb = SAMSUNG_BLUE
shape1.text = "TC 생성\n\n• 생성형 AI(LLM) 가 TC 생성\n• HLD/DLD/알고리즘 설계서 등\n• TC 생성에 과도한 시간 소요"
shape1.text_frame.paragraphs[0].font.bold = True
shape1.text_frame.paragraphs[0].font.size = Pt(16)
shape1.text_frame.paragraphs[0].font.color.rgb = WHITE

# 화살표
add_arrow(slide, Inches(4.2), Inches(2.5), Inches(1), Inches(1))

# TC 검증 박스
shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.4), Inches(2), Inches(3.5), Inches(2))
shape2.fill.solid()
shape2.fill.fore_color.rgb = SAMSUNG_LIGHT_BLUE
shape2.text = "TC 검증\n\n• 실무자가 직접 검토\n• TC 전수 검토 부담\n• 품질 판단 지표 부족"
shape2.text_frame.paragraphs[0].font.bold = True
shape2.text_frame.paragraphs[0].font.size = Pt(16)
shape2.text_frame.paragraphs[0].font.color.rgb = WHITE

# 화살표
add_arrow(slide, Inches(9.1), Inches(2.5), Inches(1), Inches(1))

# TC 수행 박스
shape3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.3), Inches(2), Inches(2.5), Inches(2))
shape3.fill.solid()
shape3.fill.fore_color.rgb = GRAY
shape3.text = "TC 수행\n\n• AI 생성 TC 품질 불안정\n• 현장/상용 문제 미반영"
shape3.text_frame.paragraphs[0].font.bold = True
shape3.text_frame.paragraphs[0].font.size = Pt(14)
shape3.text_frame.paragraphs[0].font.color.rgb = WHITE

# ==================== Slide 4: 실무자 인터뷰 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "", "실무자 인터뷰 - Pain Point")

interviews = [
    ("생성형 AI 가 3GPP 규격과 용어에 대해 학습이 부족하다", "규격/용어"),
    ("단일 FR 이 아닌 여러 FR 을 고려한 복합 Case TC 가 필요하다", "복합 Case"),
    ("FR 하나에 대해 TC 를 생성하는 데도 시간이 너무 오래 걸린다", "시간 소요"),
    ("생성해도 TC 품질이 너무 낮고 실무자의 추가적인 검토가 필수적이다", "품질/검토")
]

for i, (text, tag) in enumerate(interviews):
    top = Inches(1.5 + i * 1.5)
    bubble = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top, Inches(11.5), Inches(1.2))
    bubble.fill.solid()
    bubble.fill.fore_color.rgb = WHITE
    bubble.line.color.rgb = SAMSUNG_BLUE
    bubble.line.width = Pt(2)
    
    tf = bubble.text_frame
    tf.word_wrap = True
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = f"[{tag}] {text}"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0, 0, 0)

# ==================== Slide 5: 03. Solution Direction ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "03", "Solution Direction")

# 메인 메시지
main_msg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(11.333), Inches(2))
main_msg.fill.solid()
main_msg.fill.fore_color.rgb = SAMSUNG_BLUE
main_msg.line.color.rgb = SAMSUNG_BLUE
main_msg.text = '"AI 가 만든 Test Case, AI 가 다시 검증하고 보완하다"'
main_msg.text_frame.paragraphs[0].font.bold = True
main_msg.text_frame.paragraphs[0].font.size = Pt(28)
main_msg.text_frame.paragraphs[0].font.color.rgb = WHITE
main_msg.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# 서브 메시지
sub_msg = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(1.5))
sub_msg.text = "기존 생성 AI 를 대체하는 것이 아닌,\nRAG 기반 TC Review & Refinement System"
sub_msg.text_frame.paragraphs[0].font.size = Pt(20)
sub_msg.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
sub_msg.text_frame.paragraphs[0].font.color.rgb = SAMSUNG_LIGHT_BLUE

# ==================== Slide 6: RAG 설명 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "", "RAG(Retrieval-Augmented Generation) 란?")

# Hallucination 문제
hallu_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(5.5), Inches(2.5))
hallu_box.fill.solid()
hallu_box.fill.fore_color.rgb = RGBColor(255, 240, 240)
hallu_box.line.color.rgb = RGBColor(200, 0, 0)
hallu_box.text = "LLM Hallucination Problem\n\nAI 가 사실이 아닌 정보나 존재하지 않는 데이터를\n마치 진실인 것처럼 그럴듯하게 지어내어 답변하는 오류"
hallu_box.text_frame.paragraphs[0].font.bold = True
hallu_box.text_frame.paragraphs[0].font.size = Pt(16)

# RAG 해법
rag_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.333), Inches(1.5), Inches(5.5), Inches(2.5))
rag_box.fill.solid()
rag_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
rag_box.line.color.rgb = SAMSUNG_BLUE
rag_box.text = "RAG(Retrieval-Augmented Generation)\n\n외부 지식베이스에서 관련 정보를 검색 (Retrieval) 하여\nAI 생성 (Generation) 에 활용 → 근거 기반 답변"
rag_box.text_frame.paragraphs[0].font.bold = True
rag_box.text_frame.paragraphs[0].font.size = Pt(16)

# 비교 표
table = slide.shapes.add_table(3, 3, Inches(1), Inches(4.5), Inches(11.333), Inches(2.5))
table.table.columns[0].width = Inches(3)
table.table.columns[1].width = Inches(4)
table.table.columns[2].width = Inches(4.333)

# 헤더
headers = ["", "Fine Tuning", "RAG"]
for i, h in enumerate(headers):
    cell = table.table.cell(0, i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = SAMSUNG_BLUE
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.color.rgb = WHITE

# 내용
rows = [
    ["특징", "모델 재학습 필요", "외부 DB 검색 활용"],
    ["TC Review 적합성", "낮음 (정적)", "높음 (동적 근거 기반)"]
]
for row_idx, row in enumerate(rows):
    for col_idx, cell_text in enumerate(row):
        cell = table.table.cell(row_idx + 1, col_idx)
        cell.text = cell_text
        cell.text_frame.paragraphs[0].font.size = Pt(12)

# ==================== Slide 7: 04. System Overview ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "04", "System Overview - 3-Phase Block Workflow")

# Phase 1
p1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(3.8), Inches(2.5))
p1.fill.solid()
p1.fill.fore_color.rgb = RGBColor(230, 240, 255)
p1.line.color.rgb = SAMSUNG_BLUE
p1.text = "Phase 1: Evidence Collection\n\n• TC Metadata 추출\n• Multi-Query 검색\n• Evidence Grading\n• Evidence Pack 생성"
p1.text_frame.paragraphs[0].font.bold = True
p1.text_frame.paragraphs[0].font.size = Pt(14)

# 화살표
add_arrow(slide, Inches(4.5), Inches(2.5), Inches(0.8), Inches(0.6))

# Phase 2
p2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(1.5), Inches(3.8), Inches(2.5))
p2.fill.solid()
p2.fill.fore_color.rgb = RGBColor(200, 230, 255)
p2.line.color.rgb = SAMSUNG_LIGHT_BLUE
p2.text = "Phase 2: Validation & Refinement\n\n• 2-A: Core Validation (Req/VP/Scope)\n• 2-B: Execution (Precond/Proc/Obs/PF)\n• 2-C: Scenario & Refinement\n• Enhanced TC v2 생성"
p2.text_frame.paragraphs[0].font.bold = True
p2.text_frame.paragraphs[0].font.size = Pt(13)

# 화살표
add_arrow(slide, Inches(9.5), Inches(2.5), Inches(0.8), Inches(0.6))

# Phase 3
p3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), Inches(1.5), Inches(2.333), Inches(2.5))
p3.fill.solid()
p3.fill.fore_color.rgb = SAMSUNG_BLUE
p3.text = "Phase 3: Evaluation & Report\n\n• Quality Scoring\n• Risk Classification\n• Evidence Trace\n• Reviewer Summary"
p3.text_frame.paragraphs[0].font.bold = True
p3.text_frame.paragraphs[0].font.size = Pt(13)
p3.text_frame.paragraphs[0].font.color.rgb = WHITE

# Block 수 표시
block_info = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(1))
block_info.text = "총 28 개 Block (Phase 1: 9 개, Phase 2: 12 개, Phase 3: 7 개) - 모든 TC 가 동일한 순서로 통과"
block_info.text_frame.paragraphs[0].font.size = Pt(14)
block_info.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ==================== Slide 8: Phase 1 상세 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "", "Phase 1: Evidence Collection (9 개 Block)")

blocks_1 = [
    ("1-1", "TC Intake & Intent Understanding", "TC v1 의 목적, 관련 FR, feature, KPI 추출"),
    ("1-2", "Requirement Matching", "FR/Requirement/FD 검색"),
    ("1-3", "Design Behavior Retrieval", "HLD/DLD/Interface 문서 검색"),
    ("1-4", "Algorithm/KPI/Parameter", "알고리즘 조건, KPI, threshold 검색"),
    ("1-5", "Standard/3GPP Reference", "표준 용어, measurement 개념"),
    ("1-6", "Existing/Legacy TC", "기존 TC 와 유사 TC 검색"),
    ("1-7", "Historical Issue/Bug/Test", "과거 issue, bug, test fail 검색"),
    ("1-8", "Cross-rAPP Pattern", "다른 rAPP 의 유사 검증 pattern"),
    ("1-9", "Evidence Grading & Pack", "Strong/Supporting/Weak/Missing 등급화")
]

for i, (block_id, name, desc) in enumerate(blocks_1):
    top = Inches(1.3 + i * 0.7)
    row = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), top, Inches(12.333), Inches(0.6))
    row.fill.solid()
    row.fill.fore_color.rgb = WHITE if i % 2 == 0 else RGBColor(245, 245, 245)
    row.line.color.rgb = SAMSUNG_BLUE
    
    tf = row.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"{block_id} | {name}"
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.color.rgb = SAMSUNG_BLUE
    
    p2 = tf.add_paragraph()
    p2.text = f"    → {desc}"
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(0, 0, 0)

# ==================== Slide 9: Phase 2 상세 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "", "Phase 2: Validation & Refinement (12 개 Block)")

blocks_2a = [
    ("2-1", "Requirement Coverage", "FR 의 각 요구사항이 TC 에 반영되었는가?"),
    ("2-2", "Verification Point Coverage", "핵심 VP 가 Procedure/P-F 에 반영되었는가?"),
    ("2-3", "Scope & Target Control", "대상 rAPP/feature/NE/cell 이 명확한가?")
]

blocks_2b = [
    ("2-4", "Dependency & Precondition", "필요한 환경/데이터/설정이 명시되었는가?"),
    ("2-5", "Procedure Executability", "Test Procedure 가 실제 수행 가능한 순서인가?"),
    ("2-6", "Observability & Testability", "Expected Result 를 관측할 수 있는가?"),
    ("2-7", "Pass/Fail Clarity", "P-F Criteria 가 측정 가능한가?")
]

blocks_2c = [
    ("2-8", "Scenario Coverage", "Positive/Negative/Boundary/Exception 충분한가?"),
    ("2-9", "Historical Issue Reflection", "과거 issue/test result 가 반영되었는가?"),
    ("2-10", "Consistency & Contradiction", "TC 내부 section 간 모순이 없는가?"),
    ("2-11", "Refinement Planning", "어떤 section 을 어떤 근거로 보완할지 계획"),
    ("2-12", "Enhanced TC v2 Generation", "Refinement Plan 기반으로 TC v2 생성")
]

y = 1.3
for section, blocks in [("2-A: Core Validation", blocks_2a), ("2-B: Execution Validation", blocks_2b), ("2-C: Scenario & Refinement", blocks_2c)]:
    lbl = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(3), Inches(0.4))
    lbl.text = section
    lbl.text_frame.paragraphs[0].font.bold = True
    lbl.text_frame.paragraphs[0].font.size = Pt(12)
    lbl.text_frame.paragraphs[0].font.color.rgb = SAMSUNG_BLUE
    y += 0.4
    
    for block_id, name, desc in blocks:
        row = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(12.333), Inches(0.5))
        row.fill.solid()
        row.fill.fore_color.rgb = WHITE
        row.line.color.rgb = GRAY
        tf = row.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = f"  {block_id} {name} - {desc}"
        p.font.size = Pt(10)
        y += 0.5
    y += 0.2

# ==================== Slide 10: Phase 3 상세 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "", "Phase 3: Evaluation & Reviewer Report (7 개 Block)")

blocks_3 = [
    ("3-1", "Quality Scoring", "7 가지 차원, 100 점 만점 품질 점수"),
    ("3-2", "Risk Classification", "High/Medium/Low 3 단계 Risk 판정"),
    ("3-3", "Decision Recommendation", "Approve/Revise/Regenerate/Blocked 추천"),
    ("3-4", "Evidence Trace Table", "수정 항목과 근거 문서 연결"),
    ("3-5", "Added Scenario List", "새로 추가된 scenario 를 유형별로 정리"),
    ("3-6", "Reviewer Summary", "실무자가 먼저 볼 핵심 요약"),
    ("3-7", "Final Output Packaging", "9 가지 산출물을 고정 형식으로 패키징")
]

for i, (block_id, name, desc) in enumerate(blocks_3):
    top = Inches(1.5 + i * 0.85)
    row = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top, Inches(12.333), Inches(0.75))
    row.fill.solid()
    row.fill.fore_color.rgb = RGBColor(240, 248, 255) if i % 2 == 0 else WHITE
    row.line.color.rgb = SAMSUNG_BLUE
    
    tf = row.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"{block_id} | {name}"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = SAMSUNG_BLUE
    
    p2 = tf.add_paragraph()
    p2.text = f"    → {desc}"
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(0, 0, 0)

# ==================== Slide 11: 검증 기준 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "", "검증 기준 - 30+ 체크리스트")

categories = [
    ("Requirement Coverage", "RC-1~RC-4", "FR 의 각 요구사항이 TC 에 반영되었는가?"),
    ("Verification Point Coverage", "VP-1~VP-4", "핵심 VP 가 Procedure/P-F 에 반영되었는가?"),
    ("Scope & Target Control", "ST-1~ST-6", "대상 rAPP/feature/NE/cell 이 명확한가?"),
    ("Dependency & Precondition", "DC-1~DC-6", "필요한 환경/데이터/설정이 명시되었는가?"),
    ("Procedure Executability", "PE-1~PE-5", "Test Procedure 가 실제 수행 가능한가?"),
    ("Observability & Testability", "OT-1~OT-6", "Expected Result 를 관측할 수 있는가?"),
    ("Pass/Fail Clarity", "PF-1~PF-5", "P-F Criteria 가 측정 가능한가?"),
    ("Scenario Coverage", "SC-1~SC-8", "Positive/Negative/Boundary/Exception 충분한가?"),
    ("Historical Issue Reflection", "HI-1~HI-4", "과거 issue 가 TC 에 반영되었는가?"),
    ("Execution Feasibility", "EF-1~EF-3", "실제 수행 가능한 TC 인가?"),
    ("Template Compliance", "TC-1~TC-3", "정해진 TC format 을 따르는가?")
]

y = 1.3
for cat, ids, desc in categories:
    row = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(12.333), Inches(0.55))
    row.fill.solid()
    row.fill.fore_color.rgb = WHITE if y % 1.1 < 0.55 else RGBColor(245, 245, 245)
    row.line.color.rgb = GRAY
    
    tf = row.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"{cat} ({ids})"
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.color.rgb = SAMSUNG_BLUE
    
    p2 = tf.add_paragraph()
    p2.text = f"  {desc}"
    p2.font.size = Pt(11)
    
    y += 0.55

# ==================== Slide 12: Evidence Grading ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "", "Evidence Grading 및 Auto-Apply 규칙")

grades = [
    ("Strong", "refine_tc", "FR ID 명시적 일치\n같은 feature + 같은 KPI/threshold\n과거 issue 에서 같은 symptom + 해결 방법", "자동 보완 가능"),
    ("Supporting", "report_only", "관련 feature 언급되나 구체적 연결점 부족\n유사 issue 이나 정확한 symptom 불일치", "보고서만"),
    ("Weak", "human_review", "관련 키워드만 포함\n점수가 낮으나 완전히 관련없는 것은 아님", "Human Review 필요"),
    ("Missing", "human_review", "TC 에서 언급되었으나 검색되지 않은 핵심 요소", "Human Review 필요"),
    ("Rejected", "rejected", "관련성 낮음, 오탐지", "사용 안 함")
]

y = 1.5
for grade, usage, condition, action in grades:
    color = {"Strong": RGBColor(0, 128, 0), "Supporting": RGBColor(255, 165, 0), 
             "Weak": RGBColor(255, 200, 0), "Missing": RGBColor(255, 0, 0), "Rejected": GRAY}[grade]
    
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(12.333), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(245, 255, 245) if grade == "Strong" else WHITE
    box.line.color.rgb = color
    box.line.width = Pt(2 if grade == "Strong" else 1)
    
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"{grade} → Usage: {usage} → Action: {action}"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = color
    
    p2 = tf.add_paragraph()
    p2.text = f"  조건: {condition}"
    p2.font.size = Pt(12)
    
    y += 1.25

# ==================== Slide 13: Quality Score ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "", "Quality Score 산출 알고리즘")

# 7 가지 차원
dimensions = [
    ("Requirement Coverage", 0.25, "25%"),
    ("Scenario Completeness", 0.20, "20%"),
    ("Procedure Clarity", 0.15, "15%"),
    ("Pass/Fail Clarity", 0.15, "15%"),
    ("Data/Config Consistency", 0.10, "10%"),
    ("Evidence Traceability", 0.10, "10%"),
    ("Field Issue Awareness", 0.05, "5%")
]

# 차원 표
table = slide.shapes.add_table(8, 3, Inches(0.5), Inches(1.5), Inches(8), Inches(4.5))
table.table.columns[0].width = Inches(4)
table.table.columns[1].width = Inches(2)
table.table.columns[2].width = Inches(2)

headers = ["차원", "가중치", "기여점수"]
for i, h in enumerate(headers):
    cell = table.table.cell(0, i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = SAMSUNG_BLUE
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.color.rgb = WHITE

for i, (dim, weight, weight_str) in enumerate(dimensions):
    cell = table.table.cell(i+1, 0)
    cell.text = dim
    cell.text_frame.paragraphs[0].font.size = Pt(12)
    
    cell = table.table.cell(i+1, 1)
    cell.text = weight_str
    cell.text_frame.paragraphs[0].font.size = Pt(12)
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    cell = table.table.cell(i+1, 2)
    cell.text = f"({weight_str} × 점수)"
    cell.text_frame.paragraphs[0].font.size = Pt(11)

# 계산식
formula = slide.shapes.add_textbox(Inches(9), Inches(1.5), Inches(3.833), Inches(4))
formula.text = "Quality Score 계산식:\n\nQS = (RC × 0.25) +\n     (SC × 0.20) +\n     (PC × 0.15) +\n     (PF × 0.15) +\n     (DC × 0.10) +\n     (ET × 0.10) +\n     (FI × 0.05)\n\n총점: 100 점 만점"
formula.text_frame.paragraphs[0].font.size = Pt(13)

# ==================== Slide 14: Risk & Decision ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "", "Risk Level 판정 및 Decision Recommendation")

# Risk Matrix
risk_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(6), Inches(4.5))
risk_box.fill.solid()
risk_box.fill.fore_color.rgb = RGBColor(255, 250, 240)
risk_box.line.color.rgb = RGBColor(200, 100, 0)
risk_box.text = """Risk Level 판정 기준:

High Risk (하나라도 충족 시):
  • 핵심 requirement 누락 (Req Coverage < 60%)
  • P-F 기준 모호 (High severity)
  • 실행 불가능한 Procedure
  • 과거 high-risk issue 미반영

Medium Risk:
  • 일부 scenario 부족
  • precondition 일부 누락
  • Quality Score 60-79 점

Low Risk:
  • 사소한 표현 수정
  • Quality Score ≥ 80 점"""
risk_box.text_frame.paragraphs[0].font.bold = True
risk_box.text_frame.paragraphs[0].font.size = Pt(14)

# Decision
decision_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.5), Inches(5.833), Inches(4.5))
decision_box.fill.solid()
decision_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
decision_box.line.color.rgb = SAMSUNG_BLUE
decision_box.text = """Decision Recommendation:

✅ APPROVE (Risk=Low, Score ≥ 80)
   → 모든 Critical Issue 해결됨

⚠️ REVISE (Risk=Medium, Score ≥ 60)
   → 일부 보완 필요

❌ REGENERATE (Risk=High, Score < 60)
   → 대폭 수정 필요

🚫 BLOCKED
   → 근거 부족으로 AI 판단 불가"""
decision_box.text_frame.paragraphs[0].font.bold = True
decision_box.text_frame.paragraphs[0].font.size = Pt(14)

# ==================== Slide 15: 05. Output ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "05", "Output - 9 가지 최종 산출물")

outputs = [
    ("1. Enhanced TC v2", "보완된 Test Case (Before/After 비교 포함)"),
    ("2. Quality Report", "7 가지 차원 품질 점수, Risk Level"),
    ("3. Evidence Pack", "등급화된 근거 묶음 (Strong/Supporting/Weak/Missing)"),
    ("4. Added Scenario List", "새로 추가된 scenario 를 유형별로 정리"),
    ("5. Evidence Trace Table", "수정 항목과 근거 문서 연결"),
    ("6. Risk Indicator", "High/Medium/Low Risk 판정 결과"),
    ("7. Reviewer Summary", "실무자가 먼저 볼 핵심 요약 (Top 3 Must Check)"),
    ("8. Human Review Required List", "사람이 확인해야 할 항목"),
    ("9. Decision Recommendation", "Approve/Revise/Regenerate/Blocked 추천")
]

y = 1.5
for i, (name, desc) in enumerate(outputs):
    row = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(12.333), Inches(0.7))
    row.fill.solid()
    row.fill.fore_color.rgb = RGBColor(245, 250, 255) if i % 2 == 0 else WHITE
    row.line.color.rgb = SAMSUNG_BLUE
    
    tf = row.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = name
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = SAMSUNG_BLUE
    
    p2 = tf.add_paragraph()
    p2.text = f"  {desc}"
    p2.font.size = Pt(12)
    
    y += 0.7

# ==================== Slide 16: Before/After 예시 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "", "Before/After 예시 - TC v1 vs TC v2")

before = """[TC v1 - 기존 생성]

## Test Procedure
1. Detection Switch ON 설정
2. TA threshold 설정
3. Overshooting cell 탐지 수행
4. 결과 확인

## Pass/Fail Criteria
- COM 기능이 정상 동작한다
- TA > Threshold 인 cell 이 탐지된다"""

after = """[TC v2 - Enhanced]

## Test Procedure (Enhanced)
1. Detection Switch ON 설정 (f1=ON, f2=OFF)
2. TA threshold 설정 (thresholdTA = 120)
3. NE List 에 Indoor cell 이 포함되지 않도록 필터링
4. Overshooting cell 탐지 수행
5. f1 carrier cell 만 탐지 목록에 포함되는지 확인
6. f2 carrier cell 은 탐지 목록에 포함되지 않는지 확인

## Added Scenario
- Negative: Detection Switch OFF 상태에서 탐지되지 않음
- Boundary: TA = Threshold ±0.5dB 경계값
- Exception: CM/PM data missing 조건

## Pass/Fail Criteria (구체화)
- Detection Switch 가 ON 인 f1 carrier cell 만 detected cell list 에 포함된다
- Detection Switch 가 OFF 인 f2 carrier 의 cell 은 detected cell list 에 포함되지 않는다
- API response status 가 success 이며, detected cell count 가 expected condition 과 일치한다

## Human Review Required
- Boundary 값 (±0.5dB): 현장 허용 오차 확인 필요"""

before_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5))
before_box.fill.solid()
before_box.fill.fore_color.rgb = RGBColor(255, 245, 245)
before_box.line.color.rgb = RGBColor(200, 0, 0)
before_box.text = before
before_box.text_frame.word_wrap = True
before_box.text_frame.paragraphs[0].font.bold = True
before_box.text_frame.paragraphs[0].font.size = Pt(12)

after_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(6), Inches(5))
after_box.fill.solid()
after_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
after_box.line.color.rgb = RGBColor(0, 128, 0)
after_box.text = after
after_box.text_frame.word_wrap = True
after_box.text_frame.paragraphs[0].font.bold = True
after_box.text_frame.paragraphs[0].font.size = Pt(10)

# ==================== Slide 17: 구현 결과 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "", "구현 결과 - Quality Report & Reviewer Summary")

quality_report = """[Quality Report 예시]

## Quality Score
| 차원 | 점수 | 가중치 | 기여 |
|------|------|--------|------|
| Req Coverage | 90/100 | 25% | 22.5 |
| Scenario | 70/100 | 20% | 14.0 |
| Procedure | 80/100 | 15% | 12.0 |
| Pass/Fail | 75/100 | 15% | 11.25 |
| Data/Config | 85/100 | 10% | 8.5 |
| Evidence | 95/100 | 10% | 9.5 |
| Field Issue | 80/100 | 5% | 4.0 |
| **총점** | | | **81.75/100** |

## Decision: ✅ APPROVE
- Risk Level: Low
- 모든 Critical Issue 해결됨"""

reviewer_summary = """[Reviewer Summary 예시]

## Top 3 Must Check Items
1. Test Case #5 (Boundary) — 현장 허용 오차 확인 필요
2. Test Case #7 (Exception) — 실제 구현 여부 확인 필요
3. Pass/Fail #3 — KPI threshold 값 현행화 필요

## Major Changes
- Added Scenario: 7 개 (Negative 2, Boundary 3, Exception 2)
- Modified Procedure: 5 개 step 보완
- Revised Pass/Fail: 3 개 기준 구체화

## Human Review Required
1. Boundary 값 (±0.5dB) — 현장 검증 필요
2. Exception handling — 구현 여부 확인 필요"""

qr_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(6), Inches(5))
qr_box.fill.solid()
qr_box.fill.fore_color.rgb = WHITE
qr_box.line.color.rgb = SAMSUNG_BLUE
qr_box.text = quality_report
qr_box.text_frame.word_wrap = True
qr_box.text_frame.paragraphs[0].font.bold = True
qr_box.text_frame.paragraphs[0].font.size = Pt(11)

rs_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.5), Inches(5.833), Inches(5))
rs_box.fill.solid()
rs_box.fill.fore_color.rgb = RGBColor(245, 250, 255)
rs_box.line.color.rgb = SAMSUNG_LIGHT_BLUE
rs_box.text = reviewer_summary
rs_box.text_frame.word_wrap = True
rs_box.text_frame.paragraphs[0].font.bold = True
rs_box.text_frame.paragraphs[0].font.size = Pt(11)

# ==================== Slide 18: 06. Expected Impact ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "06", "Expected Impact - 검토 부담 감소, 품질 향상")

before_impact = """[Before - 현재]

• AI 가 기본 문맥 기반 TC 초안 생성
• 품질 판단 지표가 없어 신뢰도 낮음
• 과거 Issue 및 Edge Case 미반영
• 실무자가 TC 전체를 다시 검토 및 수정
• 검토 시간: 3~4 시간 / FR"""

after_impact = """[After - 개선]

• RAG 를 통해 TC v1 자동 검증 및 보완
• 누락 Scenario 와 모호한 P-F 자동 탐지
• 과거 Issue 기반 Enhanced TC v2 생성
• 실무자는 High Priority 영역만 집중 검토
• 검토 시간: 30 분 ~ 1 시간 / FR (70% 감소)

## 핵심 기여
1. Review Productivity: 전수검토 → Risk 중심 검토
2. Quality Accuracy: Block 기반 구조적 검증
3. Scenario Completeness: Negative/Boundary/Exception/Historical
4. Evidence Traceability: 수정 근거 추적 가능
5. Process Consistency: 동일 Block, 동일 결과
6. Extensibility: Core Rule + Domain Plugin"""

before_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5))
before_box.fill.solid()
before_box.fill.fore_color.rgb = RGBColor(255, 245, 240)
before_box.line.color.rgb = GRAY
before_box.text = before_impact
before_box.text_frame.word_wrap = True
before_box.text_frame.paragraphs[0].font.bold = True
before_box.text_frame.paragraphs[0].font.size = Pt(12)

after_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(6), Inches(5))
after_box.fill.solid()
after_box.fill.fore_color.rgb = RGBColor(240, 255, 245)
after_box.line.color.rgb = RGBColor(0, 128, 0)
after_box.text = after_impact
after_box.text_frame.word_wrap = True
after_box.text_frame.paragraphs[0].font.bold = True
after_box.text_frame.paragraphs[0].font.size = Pt(11)

# ==================== Slide 19: 07. Contribution ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "07", "Main Contribution")

contributions = [
    ("Review Productivity", "Risk-based Review, Reviewer Summary", "전수검토 부담 → Risk 중심 검토 (70% 감소)"),
    ("Quality Accuracy", "Requirement/VP/Procedure 검증", "TC 품질 신뢰도 향상 (Score 81.75/100)"),
    ("Scenario Completeness", "Negative, Boundary, Historical Issue", "누락 scenario 감소 (7 개 추가)"),
    ("Evidence Traceability", "Evidence Pack, Trace Table", "AI 수정 근거 설명 가능"),
    ("Process Consistency", "Block-based Workflow", "어떤 TC 든 동일 기준 검토 (28 개 Block)"),
    ("Extensibility", "Core Rule + Domain Plugin", "COM → 다른 rAPP 확장 가능")
]

y = 1.5
for title, keyword, effect in contributions:
    row = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(12.333), Inches(1))
    row.fill.solid()
    row.fill.fore_color.rgb = WHITE
    row.line.color.rgb = SAMSUNG_BLUE
    row.line.width = Pt(2)
    
    tf = row.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"{title}"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = SAMSUNG_BLUE
    
    p2 = tf.add_paragraph()
    p2.text = f"  키워드: {keyword}  |  효과: {effect}"
    p2.font.size = Pt(12)
    
    y += 1

# ==================== Slide 20: 08. MVP & 확장 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "08", "MVP 범위와 향후 확장 계획")

mvp = """[MVP 범위]

## 입력
• COM rAPP 관련 FR 1~3 개
• 관련 HLD/DLD 일부
• Algorithm/KPI/Parameter 일부
• 기존 생성 TC v1
• 기존 TC 예시, 과거 issue/test result 예시

## 출력
• Enhanced TC v2
• Quality Report
• Evidence Pack
• Evidence Trace Table
• Added Scenario List
• Risk Level
• Reviewer Summary
• Approve/Revise/Regenerate Recommendation

## 핵심 기능
1. TC v1 검증
2. Requirement Coverage 분석
3. Missing Scenario 탐지
4. 과거 issue 기반 추가 scenario 추천
5. TC v2 자동 보완
6. Quality Indicator 제공"""

future = """[향후 확장 계획]

## Phase 1 (완료)
• 3-Phase Block Workflow 설계
• Evidence Grading 알고리즘
• Quality Score 산출식

## Phase 2 (진행)
• Phase 1-9 Block 구현
• Phase 2-A (2-1, 2-2, 2-3) 구현
• Phase 3 Report 생성

## Phase 3 (계획)
• COM Plugin 완성
• SON Plugin 확장
• Energy Saving Plugin 확장
• 타 rAPP 적용"""

mvp_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(6), Inches(5))
mvp_box.fill.solid()
mvp_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
mvp_box.line.color.rgb = SAMSUNG_BLUE
mvp_box.text = mvp
mvp_box.text_frame.word_wrap = True
mvp_box.text_frame.paragraphs[0].font.bold = True
mvp_box.text_frame.paragraphs[0].font.size = Pt(11)

future_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.5), Inches(5.833), Inches(5))
future_box.fill.solid()
future_box.fill.fore_color.rgb = RGBColor(255, 250, 240)
future_box.line.color.rgb = RGBColor(200, 150, 0)
future_box.text = future
future_box.text_frame.word_wrap = True
future_box.text_frame.paragraphs[0].font.bold = True
future_box.text_frame.paragraphs[0].font.size = Pt(11)

# ==================== Slide 21: 향후 계획 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "", "향후 계획")

plans = [
    ("단기 (1 개월)", "• Phase 2-B/C Block 완성\n• Enhanced TC v2 자동 생성\n• End-to-End Dry-run 수행"),
    ("중기 (2 개월)", "• COM Plugin 완성\n• SON rAPP 확장\n• 실제 현장 issue 데이터 연동"),
    ("장기 (3 개월+)", "• Energy Saving rAPP 적용\n• 타 사업부 확장\n• Closed-loop TC Generation 체계 구축")
]

y = 2
for period, items in plans:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(12.333), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = SAMSUNG_BLUE
    box.line.width = Pt(2)
    
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = period
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = SAMSUNG_BLUE
    
    p2 = tf.add_paragraph()
    p2.text = items
    p2.font.size = Pt(13)
    
    y += 1.6

# ==================== Slide 22: Q&A ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Q&A
qa = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
set_title_format(qa, "Q & A")

# 감사 인사
thanks = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(1))
set_subtitle_format(thanks, "경청해 주셔서 감사합니다.")

# 발표자 정보
presenter = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11.333), Inches(0.8))
set_subtitle_format(presenter, "Samsung Networks Business RVT-RAS | 김주영")

# 저장
prs.save("output/AI_TC_Enhancement_Final_Presentation.pptx")
print("PPT 생성 완료: output/AI_TC_Enhancement_Final_Presentation.pptx")
